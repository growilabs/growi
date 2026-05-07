"""Generate minimal binary fixture files for the markitdown-extractor test suite.

Run this script once to create sample.pdf, sample.pptx, sample.xlsx, and
sample.docx in the same directory as this script.  The generated files are
committed to the repository; they do not need to be regenerated on every CI run.

Usage:
    cd services/markitdown-extractor
    uv run python tests/fixtures/generate_fixtures.py

Prerequisites (install only when regenerating fixtures):
    pip install fpdf2 python-pptx openpyxl python-docx
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def _check_available(package: str) -> bool:
    return importlib.util.find_spec(package) is not None


def generate_pdf(dest: Path) -> None:
    """Create a 3-page PDF using fpdf2."""
    if not _check_available("fpdf"):
        print(f"  [SKIP] fpdf2 not installed — cannot generate {dest.name}")
        return

    from fpdf import FPDF  # type: ignore[import]

    pdf = FPDF()
    for page_num in range(1, 4):
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, f"Page {page_num}", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", size=12)
        pdf.multi_cell(0, 8, f"This is the body text of page {page_num} in the sample PDF fixture.\n"
                              "It contains enough text to verify that the PDF extractor reads each page.")
    pdf.output(str(dest))
    print(f"  [OK]   {dest.name}")


def generate_pptx(dest: Path) -> None:
    """Create a 5-slide PPTX using python-pptx."""
    if not _check_available("pptx"):
        print(f"  [SKIP] python-pptx not installed — cannot generate {dest.name}")
        return

    from pptx import Presentation  # type: ignore[import]

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    for i in range(1, 6):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i} Title"
        tf = slide.placeholders[1].text_frame
        tf.text = f"Bullet point 1 on slide {i}"
        p = tf.add_paragraph()
        p.text = f"Bullet point 2 on slide {i}"
    prs.save(str(dest))
    print(f"  [OK]   {dest.name}")


def generate_xlsx(dest: Path) -> None:
    """Create a 2-sheet XLSX using openpyxl."""
    if not _check_available("openpyxl"):
        print(f"  [SKIP] openpyxl not installed — cannot generate {dest.name}")
        return

    import openpyxl  # type: ignore[import]

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Name", "Value"])
    ws1.append(["alpha", 1])
    ws1.append(["beta", 2])

    ws2 = wb.create_sheet("Summary")
    ws2.append(["Total", 3])
    ws2.append(["Average", 1.5])

    wb.save(str(dest))
    print(f"  [OK]   {dest.name}")


def generate_docx(dest: Path) -> None:
    """Create a single-section DOCX using python-docx."""
    if not _check_available("docx"):
        print(f"  [SKIP] python-docx not installed — cannot generate {dest.name}")
        return

    from docx import Document  # type: ignore[import]

    doc = Document()
    doc.add_heading("Sample DOCX Fixture", level=1)
    doc.add_paragraph("This is the first paragraph of the sample DOCX fixture.")
    doc.add_paragraph("This is the second paragraph, used to verify that the DOCX extractor "
                       "returns all text in a single PageInfo element.")
    doc.save(str(dest))
    print(f"  [OK]   {dest.name}")


def main() -> int:
    print("Generating binary fixture files...")
    targets = [
        (FIXTURES_DIR / "sample.pdf", generate_pdf),
        (FIXTURES_DIR / "sample.pptx", generate_pptx),
        (FIXTURES_DIR / "sample.xlsx", generate_xlsx),
        (FIXTURES_DIR / "sample.docx", generate_docx),
    ]
    for dest, generator in targets:
        if dest.exists():
            print(f"  [SKIP] {dest.name} already exists")
        else:
            generator(dest)
    print("Done.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
